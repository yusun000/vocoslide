import os
import platform
from abc import ABC, abstractmethod
from pptx import Presentation

class PPTProcessor(ABC):
    """PPT抽出の基底クラス"""
    @abstractmethod
    def extract_notes(self, pptx_path):
        pass

    @abstractmethod
    def export_slides(self, pptx_path, output_dir, pdf_path=None): # pdf_pathを追加
        pass

# --- Officeあり (Windows + PowerPoint) ---
class OfficeProcessor(PPTProcessor):
    def extract_notes(self, pptx_path):
        import win32com.client
        notes = []
        try:
            app = win32com.client.Dispatch("PowerPoint.Application")
            app.DisplayAlerts = False 
            pres = app.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
            for slide in pres.Slides:
                try:
                    note = slide.NotesPage.Shapes.Placeholders(2).TextFrame.TextRange.Text
                except:
                    note = ""
                notes.append(note)
            pres.Close()
        except Exception as e:
            print(f"Office Error: {e}")
        return notes

    def export_slides(self, pptx_path, output_dir, pdf_path=None): # 引数を合わせる
        import win32com.client
        app = win32com.client.Dispatch("PowerPoint.Application")
        pres = app.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True, WithWindow=False)
        os.makedirs(output_dir, exist_ok=True)
        for i, slide in enumerate(pres.Slides):
            slide.Export(os.path.join(os.path.abspath(output_dir), f"slide_{i+1:03d}.png"), "PNG")
        pres.Close()

# --- Officeなし (python-pptx + PDF経由) ---
class OpenProcessor(PPTProcessor):
    def extract_notes(self, pptx_path):
        from pptx import Presentation
        prs = Presentation(pptx_path)
        notes = []
        
        # スライドの並び順通りに確実に処理する
        for i, slide in enumerate(prs.slides):
            text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    text = notes_slide.notes_text_frame.text.strip()
            
            # 抽出ログを表示（デバッグ用）
            # print(f"Slide {i+1} note extracted (length: {len(text)})")
            notes.append(text)
            
        return notes

    def export_slides(self, pptx_path, output_dir, pdf_path=None): # ここを修正
        from pdf2image import convert_from_path
        # ユーザー環境のPopplerパス
        POPPLER_PATH = None

        target_pdf = pdf_path if pdf_path else pptx_path.replace(".pptx", ".pdf")
        
        if not os.path.exists(target_pdf):
            raise FileNotFoundError(f"PDF file not found: {target_pdf}")
        
        os.makedirs(output_dir, exist_ok=True)
        # 高品質設定 (dpi=150)
        images = convert_from_path(target_pdf, dpi=150, poppler_path=POPPLER_PATH)
        for i, image in enumerate(images):
            image.save(os.path.join(output_dir, f"slide_{i+1:03d}.png"), "PNG")

def get_processor():
    if platform.system() == "Windows":
        try:
            import win32com.client
            # 実際に利用可能か試みる
            app = win32com.client.Dispatch("PowerPoint.Application")
            return OfficeProcessor()
        except:
            # Officeがない場合は自動的にOpenProcessorへ
            return OpenProcessor()
    return OpenProcessor()